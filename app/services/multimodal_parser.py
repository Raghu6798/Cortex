import os
import tempfile
import base64
import mimetypes
import asyncio
import httpx
from pathlib import Path
from typing import List, Optional

from fastapi import HTTPException, Depends
from llama_parse import LlamaParse
from langchain_community.document_loaders import Docx2txtLoader, UnstructuredExcelLoader
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.messages import HumanMessage
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.config.settings import settings
from app.utils.logger import logger
from app.services.file_service import file_service

# Initialize LlamaParse
llama_parser = LlamaParse(api_key=settings.LLAMACLOUD_API_KEY, result_type="markdown")

class ParsingService:
    def __init__(self, file_service_instance):
        self.file_service = file_service_instance
        self.gemini_client = ChatGoogleGenerativeAI(
            model="gemini-1.5-flash", 
            api_key=settings.GOOGLE_API_KEY
        )

    async def parse_object(self, user_id: str, object_name: str) -> str:
        """
        Orchestrator:
        1. Generates a secure presigned URL from S3/MinIO.
        2. Downloads the file to a temp location.
        3. Routes to the correct parser based on extension.
        """
        try:
            file_url = await self.file_service.create_presigned_download_url(user_id, object_name)
        except Exception as e:
            logger.error(f"Failed to generate presigned URL: {e}")
            raise HTTPException(status_code=404, detail="File not found or access denied")

        _, ext = os.path.splitext(object_name)
        ext = ext.lower()

        try:
            if ext == ".pdf":
                return await self._process_pdf(file_url)
            elif ext in [".docx", ".doc"]:
                return await self._process_docx(file_url)
            elif ext in [".xlsx", ".xls", ".csv"]:
                return await self._process_xlsx(file_url)
            elif ext in [".pptx", ".ppt"]:
                return await self._process_pptx(file_url)
            elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
                return await self._process_image(file_url)
            else:
                return "Unsupported file format for parsing."
        except Exception as e:
            logger.error(f"Parsing failed for {object_name}: {e}", exc_info=True)
            raise HTTPException(status_code=500, detail=f"Parsing failed: {str(e)}")

    async def _download_to_temp(self, url: str, suffix: str) -> str:
        """Helper to download file from S3/MinIO URL to a temp file"""
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
            temp_path = temp_file.name
        
        async with httpx.AsyncClient() as client:
            r = await client.get(url)
            r.raise_for_status()
            Path(temp_path).write_bytes(r.content)
        
        return temp_path

    async def _process_pdf(self, url: str, pdf_path: Optional[str] = None) -> str:
        if pdf_path:
            temp_path = pdf_path
        else:
            temp_path = await self._download_to_temp(url, ".pdf")
            
        logger.info("Processing PDF with LlamaParse")
        try:
            def blocking_parse():
                return llama_parser.load_data(temp_path)
            
            documents = await asyncio.to_thread(blocking_parse)
            return "\n\n".join(doc.text for doc in documents if hasattr(doc, 'text'))
        finally:
            if not pdf_path and os.path.exists(temp_path): 
                os.remove(temp_path)

    async def _process_docx(self, url: str, docx_path: Optional[str] = None) -> str:
        if docx_path:
            temp_path = docx_path
        else:
            temp_path = await self._download_to_temp(url, ".docx")
            
        logger.info("Processing DOCX")
        try:
            return await asyncio.to_thread(
                lambda: "\n\n".join(doc.page_content for doc in Docx2txtLoader(temp_path).load())
            )
        finally:
            if not docx_path and os.path.exists(temp_path): 
                os.remove(temp_path)

    async def _process_xlsx(self, url: str, xlsx_path: Optional[str] = None) -> str:
        if xlsx_path:
            temp_path = xlsx_path
        else:
            temp_path = await self._download_to_temp(url, ".xlsx")
            
        logger.info("Processing XLSX")
        try:

            return await asyncio.to_thread(
                lambda: "\n\n".join(doc.page_content for doc in UnstructuredExcelLoader(temp_path, mode="elements").load())
            )
        finally:
            if not xlsx_path and os.path.exists(temp_path): 
                os.remove(temp_path)

    async def _process_image(self, url: str, image_path: Optional[str] = None) -> str:
        if image_path:
            temp_path = image_path
        else:
            temp_path = await self._download_to_temp(url, ".png")
            
        logger.info("Processing Image with Gemini")
        try:
            with open(temp_path, "rb") as f: 
                image_data = base64.b64encode(f.read()).decode("utf-8")
            
            mime_type, _ = mimetypes.guess_type(temp_path)
            if not mime_type:
                mime_type = "image/png"

            msg = HumanMessage(content=[
                {"type": "text", "text": "Describe this image in detail for data indexing."},
                {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_data}"}}
            ])
            
            response = await self.gemini_client.ainvoke([msg])
            return response.content
        finally:
            if not image_path and os.path.exists(temp_path): 
                os.remove(temp_path)

    async def _process_pptx(self, url: str, pptx_path: Optional[str] = None) -> str:
        if pptx_path:
            temp_path = pptx_path
        else:
            temp_path = await self._download_to_temp(url, ".pptx")
            
        logger.info("Processing PPTX (Text + Slide Images)")
        
        # We need a temp directory for extracted slide images
        temp_dir_obj = tempfile.TemporaryDirectory()
        temp_dir = temp_dir_obj.name

        try:
            prs = Presentation(temp_path)
            text_content = []
            image_tasks = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                    
       
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                        
                            image_filename = f"slide_{i}_img_{shape.shape_id}.jpg"
                            image_full_path = os.path.join(temp_dir, image_filename)
                            
                            with open(image_full_path, "wb") as f:
                                f.write(shape.image.blob)
                            
                        
                            image_tasks.append(self._describe_img(image_full_path))
                        except Exception as img_err:
                            logger.warning(f"Failed to extract image from slide {i}: {img_err}")
                
                text_content.append(f"## Slide {i+1}\n" + "\n".join(slide_text))

            
            if image_tasks:
                
                descriptions = await asyncio.gather(*image_tasks[:5])
                text_content.append("\n### Visual Content\n" + "\n".join(descriptions))

            return "\n\n".join(text_content)
            
        finally:
            temp_dir_obj.cleanup()
            if not pptx_path and os.path.exists(temp_path):
                os.remove(temp_path)

    async def _describe_img(self, path: str) -> str:
        """Helper for PPTX image description"""
        try:
            with open(path, "rb") as f: 
                b64_img = base64.b64encode(f.read()).decode()
            
            msg = HumanMessage(content=[
                {"type": "text", "text": "Summarize the visual content of this slide image."},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64_img}"}}
            ])
            res = await self.gemini_client.ainvoke([msg])
            return f"[Image Description]: {res.content}"
        except Exception as e:
            logger.error(f"Gemini image analysis failed: {e}")
            return "[Image Description Failed]"

multimodal_parser_service = ParsingService(file_service_instance=file_service)
